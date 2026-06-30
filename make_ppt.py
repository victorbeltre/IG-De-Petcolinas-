"""
Genera la presentación de propuesta de membresías PetColinas.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colores de marca ─────────────────────────────────────────────────────────
VERDE      = RGBColor(0x1a, 0x6b, 0x3a)
VERDE_DARK = RGBColor(0x14, 0x52, 0x30)
NARANJA    = RGBColor(0xd4, 0x5f, 0x1e)
DORADO     = RGBColor(0xc9, 0xa2, 0x27)
BLANCO     = RGBColor(0xFF, 0xFF, 0xFF)
CREMA      = RGBColor(0xfd, 0xf8, 0xf0)
GRIS_CLARO = RGBColor(0xf0, 0xed, 0xe8)
NEGRO      = RGBColor(0x1a, 0x1a, 0x1a)
GRIS_TEXT  = RGBColor(0x55, 0x55, 0x55)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # layout completamente en blanco

# ── Helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color, transparency=0):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if transparency:
        shape.fill.fore_color.theme_color  # ensure color set
    return shape

def add_text(slide, text, x, y, w, h, size, color, bold=False, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = "Calibri"
    return txBox

def add_logo(slide, x, y, size_inches=1.1):
    try:
        slide.shapes.add_picture("assets/logo_petcolinas.png", x, y,
                                  width=Inches(size_inches))
    except Exception:
        pass

def title_slide_bg(slide):
    """Fondo verde oscuro con círculo decorativo."""
    add_rect(slide, 0, 0, W, H, VERDE_DARK)
    # Banda naranja top
    add_rect(slide, 0, 0, W, Inches(0.12), NARANJA)
    # Banda naranja bottom
    add_rect(slide, 0, H - Inches(0.12), W, Inches(0.12), NARANJA)

def content_bg(slide, accent=True):
    """Fondo crema con banda verde arriba."""
    add_rect(slide, 0, 0, W, H, CREMA)
    if accent:
        add_rect(slide, 0, 0, W, Inches(0.12), VERDE)

def section_header(slide, color=VERDE):
    """Fondo de color sólido para slides de sección."""
    add_rect(slide, 0, 0, W, H, color)
    add_rect(slide, 0, 0, W, Inches(0.12), NARANJA)
    add_rect(slide, 0, H - Inches(0.12), W, Inches(0.12), NARANJA)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
title_slide_bg(sl)

# Chip etiqueta
chip = add_rect(sl, Inches(4.5), Inches(1.2), Inches(4.33), Inches(0.55), NARANJA)
add_text(sl, "PROPUESTA DE NEGOCIO", Inches(4.5), Inches(1.22),
         Inches(4.33), Inches(0.5), 16, BLANCO, bold=True, align=PP_ALIGN.CENTER)

# Título
add_text(sl, "Membresías PetColinas", Inches(1), Inches(2.0),
         Inches(11.33), Inches(1.2), 54, BLANCO, bold=True, align=PP_ALIGN.CENTER)

# Subtítulo dorado
add_text(sl, "Ingresos recurrentes · Clientes más leales · Crecimiento sostenido",
         Inches(1), Inches(3.3), Inches(11.33), Inches(0.7),
         22, DORADO, align=PP_ALIGN.CENTER)

# Línea decorativa
add_rect(sl, Inches(4.5), Inches(4.1), Inches(4.33), Inches(0.06), DORADO)

# Info
add_text(sl, "@petcolinas  ·  809-752-6806  ·  Plaza Las Colinas, SDO",
         Inches(1), Inches(4.4), Inches(11.33), Inches(0.5),
         16, CREMA, align=PP_ALIGN.CENTER)

add_logo(sl, Inches(5.9), Inches(5.2), size_inches=1.5)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — POR QUÉ CREAR MEMBRESÍAS
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
content_bg(sl)

add_text(sl, "¿Por qué crear membresías?", Inches(0.5), Inches(0.2),
         Inches(9), Inches(0.8), 34, VERDE, bold=True)
add_rect(sl, Inches(0.5), Inches(1.0), Inches(8), Inches(0.06), NARANJA)

razones = [
    ("💰", "Ingresos fijos garantizados",
     "Cada mes sabes exactamente cuánto vas a facturar, sin depender solo de los días buenos."),
    ("🔁", "Clientes que vuelven solos",
     "Un miembro no busca otra peluquería. Ya pagó — y viene cada 2 semanas por default."),
    ("📈", "Ticket promedio más alto",
     "Los miembros gastan más en extras y farmacia porque ya tienen la confianza y el descuento."),
    ("🐾", "Diferenciación en el mercado",
     "Pocas veterinarias en SDO tienen plan de membresía. Nos posiciona como la opción premium."),
]

y = Inches(1.2)
for icon, titulo, detalle in razones:
    add_rect(sl, Inches(0.5), y, Inches(8.5), Inches(1.1), BLANCO)
    add_rect(sl, Inches(0.5), y, Inches(0.7), Inches(1.1), VERDE)
    add_text(sl, icon, Inches(0.5), y + Inches(0.22), Inches(0.7), Inches(0.6),
             22, BLANCO, align=PP_ALIGN.CENTER)
    add_text(sl, titulo, Inches(1.3), y + Inches(0.08), Inches(5.5), Inches(0.45),
             18, NEGRO, bold=True)
    add_text(sl, detalle, Inches(1.3), y + Inches(0.52), Inches(7.2), Inches(0.5),
             13, GRIS_TEXT)
    y += Inches(1.18)

add_logo(sl, Inches(11.5), Inches(0.2), size_inches=1.2)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — NUESTROS CLIENTES (CONTEXTO)
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
content_bg(sl)

add_text(sl, "Lo que ya sabemos de nuestros clientes",
         Inches(0.5), Inches(0.2), Inches(10), Inches(0.8), 32, VERDE, bold=True)
add_rect(sl, Inches(0.5), Inches(1.0), Inches(9), Inches(0.06), NARANJA)

datos = [
    ("Frecuencia de visita",    "Cada 2–3 semanas",     "Son clientes regulares, no esporádicos."),
    ("Tamaño de mascotas",      "Pequeños y medianos",  "Poodles, Shih Tzu, Maltés, Cocker — pelo que requiere mantenimiento frecuente."),
    ("Gasto mensual actual",    "RD$1,598–RD$1,898",    "Solo en baños, sin contar extras ni vet."),
    ("Servicios que usan",      "Grooming + veterinaria","Ya combinan ambos servicios — perfectos para un plan completo."),
    ("Comportamiento esperado", "Alta retención",       "Un cliente de membresía no cambia de lugar fácilmente."),
]

y = Inches(1.2)
for i, (label, valor, nota) in enumerate(datos):
    bg = BLANCO if i % 2 == 0 else GRIS_CLARO
    add_rect(sl, Inches(0.5), y, Inches(12.2), Inches(0.9), bg)
    add_rect(sl, Inches(0.5), y, Inches(2.8), Inches(0.9), VERDE if i%2==0 else VERDE_DARK)
    add_text(sl, label, Inches(0.55), y + Inches(0.22),
             Inches(2.7), Inches(0.45), 13, BLANCO, bold=True)
    add_text(sl, valor, Inches(3.4), y + Inches(0.05),
             Inches(3), Inches(0.4), 17, NEGRO, bold=True)
    add_text(sl, nota, Inches(3.4), y + Inches(0.45),
             Inches(9.1), Inches(0.38), 12, GRIS_TEXT)
    y += Inches(1.0)

add_logo(sl, Inches(11.5), Inches(0.2), size_inches=1.2)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — PLAN BÁSICO
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
content_bg(sl)

# Header verde
add_rect(sl, 0, 0, W, Inches(1.6), VERDE)
add_rect(sl, 0, 0, W, Inches(0.12), NARANJA)
add_text(sl, "PLAN BÁSICO", Inches(0.5), Inches(0.15),
         Inches(7), Inches(0.6), 13, CREMA, bold=True)
add_text(sl, "RD$1,499", Inches(0.5), Inches(0.5),
         Inches(5), Inches(0.9), 52, BLANCO, bold=True)
add_text(sl, "/ mes · por mascota",
         Inches(3.4), Inches(0.75), Inches(3), Inches(0.5), 20, CREMA)
add_text(sl, "Para dueños que quieren grooming regular garantizado + beneficios veterinarios",
         Inches(0.5), Inches(1.2), Inches(9), Inches(0.4), 14, CREMA, italic=True)

beneficios = [
    ("✓", "2 baños completos al mes",             "Perros hasta 25 kg (pequeño o mediano)"),
    ("✓", "Corte de uñas gratis en cada visita",  "Incluido sin costo adicional"),
    ("✓", "Turno prioritario garantizado",        "Sin largas esperas — entra antes"),
    ("✓", "10% OFF en cortes y servicios extra",  "Corte completo, baño medicado, más"),
    ("✓", "10% OFF en farmacia veterinaria",       "Medicamentos y productos caninos"),
    ("✓", "10% OFF en consultas veterinarias",     "Cuando lo necesites"),
]

y = Inches(1.75)
for icon, titulo, detalle in beneficios:
    add_rect(sl, Inches(0.4), y, Inches(12.3), Inches(0.78), BLANCO)
    add_rect(sl, Inches(0.4), y, Inches(0.55), Inches(0.78), VERDE)
    add_text(sl, icon, Inches(0.4), y + Inches(0.18), Inches(0.55), Inches(0.4),
             16, BLANCO, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl, titulo, Inches(1.1), y + Inches(0.06),
             Inches(5.5), Inches(0.38), 16, NEGRO, bold=True)
    add_text(sl, detalle, Inches(1.1), y + Inches(0.43),
             Inches(11), Inches(0.3), 12, GRIS_TEXT)
    y += Inches(0.86)

# Caja de ahorro
add_rect(sl, Inches(9.5), Inches(1.75), Inches(3.3), Inches(2.2), VERDE)
add_text(sl, "AHORRO MENSUAL", Inches(9.55), Inches(1.85),
         Inches(3.2), Inches(0.4), 12, CREMA, bold=True, align=PP_ALIGN.CENTER)
add_rect(sl, Inches(9.7), Inches(2.2), Inches(2.9), Inches(0.04), DORADO)
add_text(sl, "Perro pequeño", Inches(9.55), Inches(2.35),
         Inches(3.2), Inches(0.35), 13, CREMA, align=PP_ALIGN.CENTER)
add_text(sl, "RD$299+", Inches(9.55), Inches(2.65),
         Inches(3.2), Inches(0.5), 28, DORADO, bold=True, align=PP_ALIGN.CENTER)
add_text(sl, "Perro mediano", Inches(9.55), Inches(3.1),
         Inches(3.2), Inches(0.35), 13, CREMA, align=PP_ALIGN.CENTER)
add_text(sl, "RD$599+", Inches(9.55), Inches(3.4),
         Inches(3.2), Inches(0.5), 28, DORADO, bold=True, align=PP_ALIGN.CENTER)
add_text(sl, "solo en baños · más los descuentos",
         Inches(9.55), Inches(3.85), Inches(3.2), Inches(0.3),
         10, CREMA, italic=True, align=PP_ALIGN.CENTER)

add_logo(sl, Inches(11.5), Inches(6.9), size_inches=0.9)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — PLAN PLUS
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
content_bg(sl)

# Header naranja
add_rect(sl, 0, 0, W, Inches(1.6), NARANJA)
add_rect(sl, 0, 0, W, Inches(0.12), VERDE)
add_text(sl, "PLAN PLUS  ★ PREMIUM", Inches(0.5), Inches(0.15),
         Inches(8), Inches(0.6), 13, CREMA, bold=True)
add_text(sl, "RD$2,499", Inches(0.5), Inches(0.5),
         Inches(5), Inches(0.9), 52, BLANCO, bold=True)
add_text(sl, "/ mes · por mascota",
         Inches(3.8), Inches(0.75), Inches(3), Inches(0.5), 20, CREMA)
add_text(sl, "El pack completo: grooming profesional + atención veterinaria incluida",
         Inches(0.5), Inches(1.2), Inches(9), Inches(0.4), 14, CREMA, italic=True)

beneficios_plus = [
    ("✓", "2 baños completos al mes",                    "Perros hasta 25 kg"),
    ("✓", "1 corte completo mensual",                    "A medida según la raza y preferencia"),
    ("✓", "1 consulta veterinaria de rutina mensual",    "Revisión y seguimiento — no incluye cirugías ni urgencias"),
    ("✓", "Desparasitación trimestral incluida",         "Albendazol cada 3 meses — sin costo adicional"),
    ("✓", "Turno VIP exclusivo",                         "Primer turno disponible — sin espera"),
    ("✓", "15% OFF en farmacia y servicios adicionales", "El mayor descuento disponible"),
]

y = Inches(1.75)
for icon, titulo, detalle in beneficios_plus:
    add_rect(sl, Inches(0.4), y, Inches(12.3), Inches(0.78), BLANCO)
    add_rect(sl, Inches(0.4), y, Inches(0.55), Inches(0.78), NARANJA)
    add_text(sl, icon, Inches(0.4), y + Inches(0.18), Inches(0.55), Inches(0.4),
             16, BLANCO, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl, titulo, Inches(1.1), y + Inches(0.06),
             Inches(5.5), Inches(0.38), 16, NEGRO, bold=True)
    add_text(sl, detalle, Inches(1.1), y + Inches(0.43),
             Inches(11), Inches(0.3), 12, GRIS_TEXT)
    y += Inches(0.86)

# Caja de valor
add_rect(sl, Inches(9.5), Inches(1.75), Inches(3.3), Inches(2.8), VERDE_DARK)
add_text(sl, "VALOR MENSUAL", Inches(9.55), Inches(1.85),
         Inches(3.2), Inches(0.4), 12, CREMA, bold=True, align=PP_ALIGN.CENTER)
add_rect(sl, Inches(9.7), Inches(2.2), Inches(2.9), Inches(0.04), DORADO)

items_valor = [
    ("2 baños mediano", "RD$1,898"),
    ("1 corte completo", "RD$749"),
    ("1 consulta vet", "RD$1,500"),
    ("Desparasit. (÷3)", "RD$100"),
]
y2 = Inches(2.28)
for label, val in items_valor:
    add_text(sl, label, Inches(9.55), y2, Inches(1.8), Inches(0.35), 11, CREMA)
    add_text(sl, val, Inches(11.3), y2, Inches(1.4), Inches(0.35), 11, DORADO,
             bold=True, align=PP_ALIGN.RIGHT)
    y2 += Inches(0.4)

add_rect(sl, Inches(9.7), y2 + Inches(0.0), Inches(2.9), Inches(0.04), DORADO)
add_text(sl, "Total valor", Inches(9.55), y2 + Inches(0.08),
         Inches(1.8), Inches(0.35), 12, CREMA, bold=True)
add_text(sl, "RD$4,247", Inches(11.0), y2 + Inches(0.08),
         Inches(1.7), Inches(0.35), 14, DORADO, bold=True, align=PP_ALIGN.RIGHT)
add_text(sl, "Ahorras RD$1,748", Inches(9.55), y2 + Inches(0.48),
         Inches(3.2), Inches(0.4), 15, BLANCO, bold=True, align=PP_ALIGN.CENTER)
add_text(sl, "(41% de descuento)", Inches(9.55), y2 + Inches(0.85),
         Inches(3.2), Inches(0.35), 12, CREMA, italic=True, align=PP_ALIGN.CENTER)

add_logo(sl, Inches(11.5), Inches(6.9), size_inches=0.9)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — COMPARACIÓN DE PLANES
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
content_bg(sl)

add_text(sl, "Comparación de planes", Inches(0.5), Inches(0.15),
         Inches(9), Inches(0.7), 34, VERDE, bold=True)
add_rect(sl, Inches(0.5), Inches(0.9), Inches(9), Inches(0.06), NARANJA)

col_w = Inches(3.6)
lx = Inches(3.7)
rx = Inches(7.5)

# Headers columnas
add_rect(sl, lx, Inches(1.1), col_w, Inches(0.75), VERDE)
add_text(sl, "BÁSICO  RD$1,499/mes", lx, Inches(1.18),
         col_w, Inches(0.55), 15, BLANCO, bold=True, align=PP_ALIGN.CENTER)
add_rect(sl, rx, Inches(1.1), col_w, Inches(0.75), NARANJA)
add_text(sl, "PLUS  RD$2,499/mes", rx, Inches(1.18),
         col_w, Inches(0.55), 15, BLANCO, bold=True, align=PP_ALIGN.CENTER)

filas = [
    ("2 baños mensuales",                    "✓", "✓"),
    ("Corte de uñas gratis",                 "✓", "✓"),
    ("Turno prioritario",                    "✓", "—"),
    ("Turno VIP (primer turno)",             "—", "✓"),
    ("1 corte completo mensual",             "—", "✓"),
    ("Consulta vet. de rutina mensual",      "—", "✓"),
    ("Desparasitación trimestral",           "—", "✓"),
    ("Descuento en extras y farmacia",       "10%","15%"),
    ("Descuento en consultas",               "10%","—  (incluida)"),
]

y = Inches(2.0)
for i, (servicio, bas, plu) in enumerate(filas):
    bg = BLANCO if i % 2 == 0 else GRIS_CLARO
    add_rect(sl, Inches(0.4), y, Inches(3.15), Inches(0.62), bg)
    add_rect(sl, lx, y, col_w, Inches(0.62), bg)
    add_rect(sl, rx, y, col_w, Inches(0.62), bg)
    add_text(sl, servicio, Inches(0.5), y + Inches(0.14),
             Inches(3.0), Inches(0.38), 13, NEGRO)
    col_bas = VERDE if bas == "✓" else (NARANJA if bas == "—" else NEGRO)
    col_plu = VERDE if plu == "✓" else (NARANJA if plu == "—" else NEGRO)
    add_text(sl, bas, lx, y + Inches(0.12), col_w, Inches(0.38),
             17, col_bas, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl, plu, rx, y + Inches(0.12), col_w, Inches(0.38),
             17, col_plu, bold=True, align=PP_ALIGN.CENTER)
    y += Inches(0.66)

add_logo(sl, Inches(11.5), Inches(0.15), size_inches=1.1)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — OFERTA SEGUNDA MASCOTA
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
content_bg(sl)

add_text(sl, "Oferta Segunda Mascota", Inches(0.5), Inches(0.15),
         Inches(9), Inches(0.7), 34, VERDE, bold=True)
add_rect(sl, Inches(0.5), Inches(0.9), Inches(10), Inches(0.06), NARANJA)

add_text(sl, "50% de descuento en la membresía del mismo plan para una segunda mascota.",
         Inches(0.5), Inches(1.05), Inches(11), Inches(0.55),
         16, GRIS_TEXT, italic=True)

for i, (plan, precio1, precio2, total, ahorro) in enumerate([
    ("Plan BÁSICO",  "RD$1,499", "RD$750",   "RD$2,249/mes", "vs RD$2,998 pagando individual"),
    ("Plan PLUS",    "RD$2,499", "RD$1,250", "RD$3,749/mes", "vs RD$4,998 pagando individual"),
]):
    bx = Inches(0.4 + i * 6.5)
    add_rect(sl, bx, Inches(1.75), Inches(6.0), Inches(4.8),
             VERDE if i == 0 else NARANJA)
    add_text(sl, plan, bx, Inches(1.9), Inches(6.0), Inches(0.55),
             20, BLANCO, bold=True, align=PP_ALIGN.CENTER)
    add_rect(sl, bx + Inches(0.5), Inches(2.5), Inches(5.0), Inches(0.05), BLANCO)

    add_text(sl, "1ra mascota", bx, Inches(2.65), Inches(2.9), Inches(0.45), 14, CREMA, align=PP_ALIGN.CENTER)
    add_text(sl, precio1, bx + Inches(3.0), Inches(2.65), Inches(2.9), Inches(0.45),
             18, DORADO, bold=True, align=PP_ALIGN.CENTER)

    add_text(sl, "2da mascota", bx, Inches(3.15), Inches(2.9), Inches(0.45), 14, CREMA, align=PP_ALIGN.CENTER)
    add_text(sl, precio2, bx + Inches(3.0), Inches(3.15), Inches(2.9), Inches(0.45),
             18, DORADO, bold=True, align=PP_ALIGN.CENTER)

    add_rect(sl, bx + Inches(0.5), Inches(3.75), Inches(5.0), Inches(0.05), DORADO)
    add_text(sl, "TOTAL MENSUAL", bx, Inches(3.88), Inches(6.0), Inches(0.4),
             12, CREMA, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl, total, bx, Inches(4.28), Inches(6.0), Inches(0.75),
             32, BLANCO, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl, ahorro, bx, Inches(5.05), Inches(6.0), Inches(0.4),
             12, CREMA, italic=True, align=PP_ALIGN.CENTER)

add_text(sl, "★  Esta oferta aplica solo para miembros activos y para mascotas del mismo dueño.",
         Inches(0.4), Inches(6.8), Inches(12), Inches(0.45),
         13, GRIS_TEXT, italic=True)
add_logo(sl, Inches(11.5), Inches(0.15), size_inches=1.1)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — PROYECCIÓN DE INGRESOS
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
content_bg(sl)

add_text(sl, "¿Qué significa esto para el negocio?",
         Inches(0.5), Inches(0.15), Inches(11), Inches(0.7), 32, VERDE, bold=True)
add_rect(sl, Inches(0.5), Inches(0.9), Inches(11), Inches(0.06), NARANJA)

add_text(sl, "Proyección de ingresos mensuales recurrentes — solo por membresías",
         Inches(0.5), Inches(1.05), Inches(11), Inches(0.4), 14, GRIS_TEXT, italic=True)

escenarios = [
    ("Inicio conservador",  "10 Básico + 5 Plus",   "RD$27,490/mes",  "RD$329,880/año"),
    ("Crecimiento moderado","20 Básico + 15 Plus",  "RD$67,460/mes",  "RD$809,520/año"),
    ("Escenario ideal",     "40 Básico + 30 Plus",  "RD$134,920/mes", "RD$1,619,040/año"),
]

y = Inches(1.65)
for i, (escenario, mix, mensual, anual) in enumerate(escenarios):
    colors = [VERDE, NARANJA, VERDE_DARK]
    add_rect(sl, Inches(0.4), y, Inches(12.2), Inches(1.2), colors[i])
    add_text(sl, escenario.upper(), Inches(0.6), y + Inches(0.08),
             Inches(2.8), Inches(0.45), 13, CREMA, bold=True)
    add_text(sl, mix, Inches(0.6), y + Inches(0.55),
             Inches(2.8), Inches(0.45), 12, CREMA, italic=True)
    add_text(sl, mensual, Inches(4.5), y + Inches(0.22),
             Inches(3.8), Inches(0.7), 30, BLANCO, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl, "mensuales", Inches(4.5), y + Inches(0.82),
             Inches(3.8), Inches(0.3), 11, CREMA, align=PP_ALIGN.CENTER)
    add_text(sl, anual, Inches(8.5), y + Inches(0.22),
             Inches(3.8), Inches(0.7), 30, DORADO, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl, "al año", Inches(8.5), y + Inches(0.82),
             Inches(3.8), Inches(0.3), 11, CREMA, align=PP_ALIGN.CENTER)
    y += Inches(1.3)

add_text(sl, "★  Ingresos fijos antes de contar los servicios individuales, los extras de los miembros y la farmacia.",
         Inches(0.4), Inches(6.8), Inches(12), Inches(0.45), 12, GRIS_TEXT, italic=True)
add_logo(sl, Inches(11.5), Inches(0.15), size_inches=1.1)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — PASOS PARA EL LANZAMIENTO
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
content_bg(sl)

add_text(sl, "Plan de lanzamiento", Inches(0.5), Inches(0.15),
         Inches(9), Inches(0.7), 34, VERDE, bold=True)
add_rect(sl, Inches(0.5), Inches(0.9), Inches(9), Inches(0.06), NARANJA)

pasos = [
    ("1", "Semana 1",   "Aprobación interna",
     "Ajustar los detalles del plan, imprimir materiales, capacitar al equipo."),
    ("2", "Semana 2",   "Pre-lanzamiento en redes",
     "Publicar teasers en Instagram con conteo regresivo — generar expectativa."),
    ("3", "Semana 3",   "Lanzamiento oficial",
     "Post de anuncio completo en Instagram + oferta de 'fundadores' para los primeros 10 miembros."),
    ("4", "Semana 4+",  "Captación activa",
     "Ofrecer la membresía a cada cliente que venga, medir conversión y ajustar."),
]

y = Inches(1.15)
for num, semana, titulo, detalle in pasos:
    add_rect(sl, Inches(0.4), y, Inches(12.2), Inches(1.3), BLANCO if int(num)%2==1 else GRIS_CLARO)
    add_rect(sl, Inches(0.4), y, Inches(0.75), Inches(1.3), VERDE if int(num)%2==1 else NARANJA)
    add_text(sl, num, Inches(0.4), y + Inches(0.35), Inches(0.75), Inches(0.6),
             26, BLANCO, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl, semana, Inches(1.3), y + Inches(0.08), Inches(2), Inches(0.38), 12, GRIS_TEXT, bold=True)
    add_text(sl, titulo, Inches(1.3), y + Inches(0.4), Inches(4.5), Inches(0.45), 17, NEGRO, bold=True)
    add_text(sl, detalle, Inches(5.9), y + Inches(0.32), Inches(6.6), Inches(0.7), 13, GRIS_TEXT)
    y += Inches(1.42)

add_logo(sl, Inches(11.5), Inches(6.9), size_inches=0.9)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CIERRE
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
title_slide_bg(sl)

add_text(sl, "Las membresías no son un gasto para el cliente.", Inches(1), Inches(1.6),
         Inches(11.33), Inches(0.9), 38, BLANCO, bold=True, align=PP_ALIGN.CENTER)
add_text(sl, "Son la razón para que siempre vuelva a nosotros.",
         Inches(1), Inches(2.55), Inches(11.33), Inches(0.9),
         38, DORADO, bold=True, align=PP_ALIGN.CENTER)

add_rect(sl, Inches(3.5), Inches(3.55), Inches(6.33), Inches(0.07), BLANCO)

add_text(sl, "Con solo 20 miembros cubrimos nómina.",
         Inches(1), Inches(3.8), Inches(11.33), Inches(0.5), 20, CREMA, align=PP_ALIGN.CENTER)
add_text(sl, "Con 50 miembros construimos el negocio.",
         Inches(1), Inches(4.25), Inches(11.33), Inches(0.5), 20, CREMA, align=PP_ALIGN.CENTER)

add_logo(sl, Inches(5.6), Inches(5.1), size_inches=2.1)

# ── Guardar ──────────────────────────────────────────────────────────────────
out = "PetColinas_Membresias_Propuesta.pptx"
prs.save(out)
print(f"Presentacion guardada: {out}")
